"""Author 3 .docx + 3 .pptx source training materials across 6 distinct domains.

Every doc has SPECIFIC, TRACEABLE FACTS (company names, dollar amounts, named
people, concrete thresholds). This lets a fidelity agent later verify whether
the LMS-generated course stays grounded or drifts into made-up content.

Output: /tmp/grounding_test/*.docx, *.pptx
"""
import os
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt

OUT = "/tmp/grounding_test"
os.makedirs(OUT, exist_ok=True)


# ============================================================
# DOCX 1: Internal SOP for claim triage at MeridianInsure (fictional)
# ============================================================
def docx_meridian_claim_triage():
    d = Document()
    d.add_heading("MeridianInsure Claim Triage SOP v4.2", 0)
    d.add_paragraph("Effective date: March 14, 2026. Owner: Priya Rao, Head of Claims Ops. Approved by: CFO Daniel Kwon.")

    d.add_heading("1. Scope and audience", 1)
    d.add_paragraph(
        "This SOP applies to all Tier-1 claim adjusters at MeridianInsure (Phoenix and Newark "
        "operations centers) handling auto and homeowner claims between $2,500 and $75,000. "
        "Claims below $2,500 follow the auto-settle pipeline (SOP v4.1, Section 3) and claims "
        "above $75,000 are escalated to Complex Claims Unit led by Marcus Delgado."
    )

    d.add_heading("2. The MERIDIAN checklist (seven stages)", 1)
    stages = [
        ("M — Match", "Within 30 minutes of claim intake, match the claim to the policyholder's active policy using claim_id format CLM-YYYY-NNNNNN. Flag mismatches in Salesforce queue 'CLM-NOMATCH'."),
        ("E — Evidence", "Collect at least 4 photos, 1 police report (if collision), and 1 first-notice-of-loss call transcript. Evidence completeness threshold: 80% or claim cannot proceed to Appraise."),
        ("R — Reach-out", "Schedule first policyholder call within 24 hours. Use Zendesk macro 'First-touch-auto-v4'. If no response after 3 attempts over 72 hours, mark 'UNRESPONSIVE' and escalate to Rahul Srinivasan's subrogation team."),
        ("I — Investigate", "Run fraud score via LexisNexis API. Scores above 0.74 route to Special Investigations Unit (SIU). Document all red flags with timestamp."),
        ("D — Determine", "Calculate loss-adjusted cost using the LACP-2026 rate card (updated quarterly by Controller Amelia Song). Apply depreciation curve per asset class."),
        ("I — Issue", "Generate settlement offer letter using Template ID LTR-4471-v3. Must include exactly three things: covered amount, deductible applied, appeal process (30 days from postmark)."),
        ("A — Aftercare", "30-day follow-up NPS survey. Scores below 7 trigger a QA review by Priya Rao's team. QA findings feed into monthly calibration meeting held every 2nd Thursday at 2pm PST."),
    ]
    for name, body in stages:
        d.add_heading(name, 2)
        d.add_paragraph(body)

    d.add_heading("3. Escalation thresholds", 1)
    d.add_paragraph(
        "Any of the following triggers immediate escalation to Karen Liu (VP Claims) within 1 business hour: "
        "(a) claim value exceeds $150,000; (b) fraud score above 0.91; (c) policyholder is a minor or requires ADA accommodations; "
        "(d) litigation threat language detected by Claude classifier; (e) claim involves bodily injury requiring medical specialist review."
    )

    d.add_heading("4. Known pitfalls", 1)
    d.add_paragraph(
        "Adjusters frequently skip the fraud-score step (Stage I) when policyholders are long-tenured (10+ years). "
        "This created a $2.3M loss in Q3 2025 (the 'Oakridge cluster' — 14 duplicate claims from the same repair shop). "
        "Do not skip the fraud step regardless of tenure."
    )
    d.add_paragraph(
        "Second common failure: the 72-hour UNRESPONSIVE marker is often left past 7 days, which blocks subrogation. "
        "If no policyholder response in 72 hours, mark immediately."
    )

    d.add_heading("5. KPIs tracked monthly", 1)
    kpis = [
        "Cycle time from intake to settlement: target <= 8.5 business days",
        "Fraud catch rate: target >= 11.2% (industry benchmark 9.8%)",
        "Policyholder NPS: target >= 64 (Q4 2025 was 59)",
        "Subrogation recovery: target >= $1.4M/month (Q4 2025 was $1.18M)",
        "SIU referral accuracy: target >= 82% true-positive rate",
    ]
    for k in kpis:
        d.add_paragraph(k, style="List Bullet")

    path = f"{OUT}/meridian_claim_triage_sop.docx"
    d.save(path)
    return path


# ============================================================
# DOCX 2: Onboarding handbook for Junior Data Engineers at Kelvingrove Analytics (fictional)
# ============================================================
def docx_kelvingrove_onboarding():
    d = Document()
    d.add_heading("Kelvingrove Analytics — Junior Data Engineer Onboarding", 0)
    d.add_paragraph("Last updated: January 22, 2026. Owner: Evelyn Harsha (Director of Data Platform).")

    d.add_heading("Week 1: The stack and your first PR", 1)
    d.add_paragraph(
        "You will be ramping into our data platform: Snowflake (warehouse region us-east-2), dbt Cloud, "
        "Airflow 2.7 deployed on Astronomer, Looker for BI, Segment for event capture. "
        "Your first commit: fix the `stg_hubspot__contacts` model where email_domain is NULL for 4.2% of rows — "
        "root cause is a regex boundary bug in staging. Target: merge by end of Week 1."
    )

    d.add_heading("Week 2: On-call shadowing", 1)
    d.add_paragraph(
        "Shadow Arjun Kapoor's on-call rotation. Pay attention to the Airflow SLA miss in our "
        "`revenue_recognition_daily` DAG — it fires at 04:00 UTC and must complete before 06:30 UTC "
        "for Controller Amelia Song's manual close. Expected failure modes: (1) Stripe API rate limit "
        "(retry logic in `task_pull_invoices`), (2) Snowflake warehouse cold-start on WH_XSMALL_ANALYST, "
        "(3) dbt test failure on `fct_revenue.amount_cents` non-null assertion."
    )

    d.add_heading("Week 3-4: Your first project — the churn-risk rewrite", 1)
    d.add_paragraph(
        "Our current churn-risk model was authored by ex-employee Ramesh K in 2023. It has been "
        "producing a 47% false-positive rate which is hurting the CS team's credibility. Your job: "
        "rewrite as a dbt model using these features: logins_per_week_last_4, support_tickets_open, "
        "days_since_last_feature_adoption, MRR_class (NEW/EXPAND/STEADY/DECLINE). Target FPR: <= 18%."
    )

    d.add_heading("Expectations", 1)
    d.add_paragraph(
        "By Day 30: you can ship a dbt model PR without Arjun's review. "
        "By Day 60: you can handle P2 on-call alerts autonomously (but still escalate P1 to Arjun). "
        "By Day 90: you own the churn-risk model and ship the rewrite. "
        "Performance review with Evelyn Harsha on Day 90."
    )

    d.add_heading("Tooling you'll set up Day 1", 1)
    d.add_paragraph(
        "Install: dbt Cloud (profile 'kelvingrove_dev'), Snowflake CLI (warehouse WH_DEV_XSMALL), "
        "Git SSH key, VS Code with our Kelvingrove extension pack (dbt-power-user, Snowflake-vsc, "
        "and our internal 'kelv-lint' for SQL style). Ask Priyanka S in Platform-Ops if any of the "
        "above fails to install — do NOT try to debug yourself. Platform-Ops office hours: "
        "Tuesdays and Thursdays 2-3pm PST in Zoom room 'platform-office-hours'."
    )

    path = f"{OUT}/kelvingrove_onboarding.docx"
    d.save(path)
    return path


# ============================================================
# DOCX 3: Compliance bulletin for GDPR DSAR handling at Zephyr Health (fictional)
# ============================================================
def docx_zephyr_dsar_compliance():
    d = Document()
    d.add_heading("Zephyr Health — GDPR Data-Subject Access Request (DSAR) Handling", 0)
    d.add_paragraph("Bulletin 2026-03. Owner: Rivka Mendelsohn (Data Protection Officer). Effective: April 1, 2026.")

    d.add_heading("Regulatory context", 1)
    d.add_paragraph(
        "Under Article 15 of GDPR, EU data subjects may request: (a) confirmation their data is being processed, "
        "(b) a copy of the data, (c) the purposes of processing, (d) the categories of data, (e) the recipients "
        "of data, (f) retention period, (g) right to rectification/erasure. Response deadline: 30 calendar days "
        "from intake (extensible to 90 days for complex requests with justified notification to the subject)."
    )

    d.add_heading("Zephyr intake and routing", 1)
    d.add_paragraph(
        "All DSAR requests enter via dsar@zephyr.health (monitored by Rivka's team). Intake triage within 2 business days: "
        "(1) verify identity via our IdP partner (Persona), reject requests that fail identity match twice; "
        "(2) scope the request — full-history, last-12-months, or specific categories; "
        "(3) assign ticket in Jira project 'DSAR' with label 'intake' and SLA countdown. "
        "Do NOT begin data-extraction until identity is verified — this was violated in the 2024 'Freeman incident' "
        "where a $340,000 regulatory fine was avoided only because the data hadn't left our environment."
    )

    d.add_heading("Systems we must query per DSAR", 1)
    systems = [
        "Clinical records in Epic (via our Epic-sync warehouse in Snowflake)",
        "Billing records in NetSuite ERP",
        "Customer support records in Zendesk (retention: 7 years)",
        "Marketing records in HubSpot (retention: per user's consent grant)",
        "Event analytics in Amplitude (retention: 3 years rolling)",
        "Third-party processor records — we maintain a current sub-processor list at zephyr.health/privacy/subprocessors",
        "Physical archive at Iron Mountain facility IM-PX-0482 (Phoenix, AZ) — allow 5 business days for retrieval",
    ]
    for s in systems:
        d.add_paragraph(s, style="List Bullet")

    d.add_heading("Response packaging standard", 1)
    d.add_paragraph(
        "Deliver via encrypted portal (Tresorit). Package contents: (a) structured data export in JSON per Article 20 portability spec; "
        "(b) cover letter on Zephyr letterhead signed by DPO Rivka Mendelsohn; "
        "(c) redactions applied for any data containing a third party's PII (Article 15 paragraph 4 limitation); "
        "(d) a plain-English summary not exceeding 600 words; "
        "(e) a 'your rights' appendix listing erasure, rectification, restriction, objection, portability, and right to lodge complaint with supervisory authority."
    )

    d.add_heading("Escalation", 1)
    d.add_paragraph(
        "Escalate to General Counsel David Park within 24 hours if: (1) the subject references litigation, "
        "(2) the data includes protected health information requiring HIPAA authorization review, "
        "(3) the subject is a minor (custodial parent consent required), "
        "(4) the request involves a data breach we have not yet disclosed, "
        "(5) the subject is in a jurisdiction outside EU/UK/California requiring separate legal analysis."
    )

    path = f"{OUT}/zephyr_dsar_compliance.docx"
    d.save(path)
    return path


# ============================================================
# PPTX 1: Sales playbook — Pre-Series-B Enterprise Upsell at Orbital-Nova (fictional)
# ============================================================
def pptx_orbital_enterprise_upsell():
    p = Presentation()
    p.slide_width = Inches(13.33)
    p.slide_height = Inches(7.5)
    blank = p.slide_layouts[6]

    def slide(title, bullets):
        s = p.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8)).text_frame
        tb.text = title
        tb.paragraphs[0].runs[0].font.size = Pt(28)
        tb.paragraphs[0].runs[0].font.bold = True
        body = s.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8)).text_frame
        for b in bullets:
            para = body.add_paragraph()
            para.text = b
            para.runs[0].font.size = Pt(18)

    slide("Orbital-Nova — Enterprise Upsell Playbook v2.1 (Q2 2026)",
        ["Owner: Kenji Yamamoto, VP Revenue",
         "Target customer segment: Series B/C companies, $50M-$500M ARR, 200-2000 employees",
         "Our product: the Orbital Platform — observability + incident management",
         "Upsell motion: move existing $36K/yr Professional customers to $120K/yr Enterprise tier",
         "Why Enterprise: SSO, audit log export, dedicated Slack support, 99.99% SLA, compliance attestations"])

    slide("Who to call — the ICP Scorecard",
        ["Company size 200-2000 employees (must)",
         "Current spend $36K-$48K/yr (sweet spot for 3x expand)",
         "Has >=3 SRE engineers on staff (LinkedIn check)",
         "Uses Okta or Azure AD (SSO trigger)",
         "Has a compliance-sensitive industry (healthcare/fintech/govtech) — triggers audit-log upgrade",
         "Champion = usually a Staff+ SRE or VP Eng",
         "Economic buyer = CTO or VP Eng"])

    slide("The Seven-Step Orbital Motion",
        ["1. Usage-signal detection — Usage-ops (Amaya) flags accounts with >=40 active users + alert volume >100/day",
         "2. Champion outreach — use template BDR-EXPAND-v4, 3-touch sequence over 10 days",
         "3. Discovery call (30 min) — questions script at go/orbital-discovery-qs",
         "4. Technical deep-dive (45 min) — SE Priya Naresh demos audit-log + SSO + on-prem agent",
         "5. Security questionnaire — 47-item template, we auto-fill via Vanta integration",
         "6. Procurement — route to our Legal (Delia Jordan-Smith) for MSA negotiation",
         "7. Close — offer 10% discount for annual prepay, 5% for quarterly"])

    slide("Common objections and our positions",
        ["'Enterprise is 3x the price, what's the ROI?' — cite the $420K MTTR savings paper (case-study-CloudSpire-2025)",
         "'We already use Datadog' — Orbital differentiator: we ingest their OTEL output, we don't replace their APM",
         "'Our compliance team needs SOC 2 Type II' — we have this; Vanta Trust Center link is vanta.trust.orbital-nova.com",
         "'Can we self-host?' — NO. Only Enterprise Plus ($240K/yr) has self-host option, and only for healthcare/govtech",
         "'Need proof-of-value first' — offer a 45-day POV with signed mutual NDA, scoped to ONE service"])

    slide("The three Enterprise-tier features that close deals",
        ["SAML SSO via Okta or Azure AD — 67% of our Enterprise deals cite this as must-have",
         "Audit log export (S3 / GCS / Splunk) — compliance trigger, especially financial services",
         "Dedicated private Slack with our VP Customer Success Daniela Falke — 'executive support'"])

    slide("Red flags — walk away signals",
        ["Champion won't intro us to economic buyer after 2 asks — deal is dead",
         "Company just laid off >10% of engineering team — freeze, revisit in 6 months",
         "Competitor (Datadog, New Relic, Splunk) is on an overlapping renewal — negotiation dynamics bad",
         "Our support has >3 open P1 tickets with them unresolved >30 days — do not expand before resolved",
         "Champion is leaving (LinkedIn 'open to work') — pause until replacement identified"])

    slide("Metrics Kenji tracks weekly",
        ["Pipeline coverage: 3.2x quota (target >=3.0)",
         "Win rate on expand deals: target >=42%",
         "Average deal size: target $138K (up from $112K in Q4 2025)",
         "Time to close from first call: target <= 65 days",
         "Discount authority: AEs can go to 12%, anything lower needs Kenji approval"])

    path = f"{OUT}/orbital_enterprise_upsell.pptx"
    p.save(path)
    return path


# ============================================================
# PPTX 2: Clinical research protocol for Trial-047 at Lumen-Bio (fictional)
# ============================================================
def pptx_lumen_trial047():
    p = Presentation()
    p.slide_width = Inches(13.33); p.slide_height = Inches(7.5)
    blank = p.slide_layouts[6]
    def slide(t, bs):
        s = p.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8)).text_frame
        tb.text = t; tb.paragraphs[0].runs[0].font.size = Pt(26); tb.paragraphs[0].runs[0].font.bold = True
        body = s.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8)).text_frame
        for b in bs:
            para = body.add_paragraph(); para.text = b; para.runs[0].font.size = Pt(17)

    slide("Lumen-Bio Trial-047: Phase 2 Protocol Training",
        ["Molecule: LB-4721 (oral kinase inhibitor, class JAK2-selective)",
         "Indication: moderate plaque psoriasis, PASI >=12",
         "PI: Dr. Meena Bhattacharya, MD, PhD (San Francisco site)",
         "Sponsor regulatory lead: Graziano Pellini (FDA interactions)",
         "Study arms: 80 patients, 2:1 randomization (treatment vs placebo)",
         "Primary endpoint: PASI-75 at Week 16",
         "Secondary: DLQI improvement >=5, IGA 0/1"])

    slide("Inclusion criteria",
        ["Age 18-70, any sex",
         "Confirmed plaque psoriasis diagnosis by board-certified dermatologist",
         "Baseline PASI >=12 AND BSA >=10%",
         "Must have failed at least one prior systemic therapy (biologic OR oral)",
         "Negative pregnancy test at screening (for patients of childbearing potential)",
         "Able to swallow oral capsules",
         "Signed informed consent, form version ICF-047-v3 dated 2026-02-18"])

    slide("Exclusion criteria",
        ["Guttate, erythrodermic, or pustular psoriasis subtypes",
         "Prior JAK inhibitor use within last 12 weeks",
         "Active or chronic infection requiring systemic therapy",
         "History of malignancy within last 5 years (except non-melanoma skin cancer)",
         "Pregnancy or breastfeeding",
         "eGFR < 60 mL/min/1.73m^2",
         "ALT or AST > 2x upper limit of normal"])

    slide("Dosing schedule",
        ["Week 0 (baseline): randomize, first dose in clinic, 30-min post-dose observation",
         "Weeks 1, 2, 4, 8, 12, 16: in-clinic visits for PASI and safety bloodwork",
         "Daily dosing: LB-4721 15mg twice daily OR matched placebo twice daily",
         "Dose reduction allowed to 15mg once daily per investigator judgment for grade 2+ adverse events",
         "Discontinuation for grade 3+ hepatic or hematologic adverse events"])

    slide("Adverse event reporting",
        ["Any serious adverse event (SAE) must be reported to Lumen-Bio medical monitor within 24 hours",
         "Monitor: Dr. Vinod Kumaran, reachable via the 24/7 hotline +1-415-555-0194",
         "Use MedDRA coding on all AE reports; use CTCAE v5.0 grading for severity",
         "Mandatory FDA MedWatch submission for all SAEs involving hospitalization or death",
         "Site PI must sign off on all SAE narratives within 72 hours"])

    slide("Site monitoring visits",
        ["Pre-study: site qualification visit by Kimberly Oduya (CRA)",
         "Initiation: 2-day on-site training, 2026-04-02 and 2026-04-03",
         "Routine monitoring: every 6 weeks during active enrollment",
         "Close-out: full SDV (source data verification) of 100% of primary endpoint data",
         "Site payment schedule: $4,200 per enrolled patient + $8,500 site-startup fee",
         "All queries resolved in Medidata Rave within 10 business days"])

    path = f"{OUT}/lumen_trial_047_protocol.pptx"
    p.save(path)
    return path


# ============================================================
# PPTX 3: Crisis comms playbook for Halcyon Cybersecurity (fictional)
# ============================================================
def pptx_halcyon_crisis_comms():
    p = Presentation()
    p.slide_width = Inches(13.33); p.slide_height = Inches(7.5)
    blank = p.slide_layouts[6]
    def slide(t, bs):
        s = p.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8)).text_frame
        tb.text = t; tb.paragraphs[0].runs[0].font.size = Pt(26); tb.paragraphs[0].runs[0].font.bold = True
        body = s.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8)).text_frame
        for b in bs:
            para = body.add_paragraph(); para.text = b; para.runs[0].font.size = Pt(17)

    slide("Halcyon Cybersecurity — Breach Crisis Comms Playbook",
        ["Version: 2.3, last updated March 1, 2026",
         "Owner: Sonia Amador (Chief Communications Officer)",
         "Distribution: C-suite, Legal (GC: Marcus Hlongwane), Customer Success leaders, Board observer",
         "Purpose: standardize external communications in the first 96 hours of a confirmed security breach"])

    slide("Activation — who calls it",
        ["CEO Ariel Nazarova activates the playbook within 4 hours of a CONFIRMED breach",
         "CISO Elena Kowalski confirms technical details with the IR lead (Diego Reyes)",
         "CCO Sonia Amador owns all external comms; Legal Marcus Hlongwane signs off on each public statement",
         "Board chair Harold Lindgren is notified via private line within 1 hour of activation"])

    slide("Hour 0-4: Internal mobilization",
        ["Spin up 'war room' in conference room H-211 (Austin HQ) OR Zoom 'halcyon-crisis-comms'",
         "First communication goes to ALL employees via Slack #all-hands: 'We are investigating a security incident. Do not discuss externally until we issue guidance.'",
         "Deactivate all employee social-media auto-posting tools (Hootsuite, Buffer) until further notice",
         "Prepare a holding statement draft (CCO template 'halcyon-holding-v4') within 2 hours",
         "GC Marcus reviews all statements before release"])

    slide("Hour 4-24: First external communications",
        ["Customer notification: individualized email from Sonia Amador to top-50 enterprise accounts",
         "Public holding statement on halcyon.security/security-update (not press release yet)",
         "Regulatory notifications: CISA within 24h if critical infrastructure; state AG per CCPA/state breach laws",
         "Do NOT confirm scope until IR has finalized blast-radius assessment",
         "Media hold — route all inquiries to press@halcyon.security; Sonia signs off on any comment",
         "Customer Support: spin up '#breach-support' Slack channel, arm CS leads with Q&A doc (template 'halcyon-cs-qa-v3')"])

    slide("Hour 24-72: Detailed disclosure",
        ["Detailed customer notification: scope, timeline, what was/was not accessed, remediation steps",
         "Public-facing detailed post from CEO Ariel Nazarova — reviewed by Marcus Hlongwane AND Harold Lindgren before publish",
         "Press release via Business Wire (account HLCN-25582)",
         "Offer affected customers: free 24 months of credit monitoring via Kroll (contract ID KRL-HCN-22)",
         "File Form 8-K disclosure if materially significant (Legal + CFO Joel Weidenfeld)",
         "Post to halcyon.security/transparency within 72 hours with timeline and RCA narrative"])

    slide("Tone guidelines",
        ["Empathetic, not defensive. Example: 'We are deeply sorry' not 'We regret any inconvenience'.",
         "Concrete, not vague. Example: '~14,000 customer records' not 'a limited number of records'.",
         "No blame externalization — do NOT attribute to a third-party vendor until IR is certain",
         "Every public statement must answer: What happened, when did we detect, what have we done, what should you do",
         "Avoid 'hackers' — prefer 'unauthorized party' or 'threat actor' (legal preference)"])

    slide("Post-incident (day 7+)",
        ["30-day post-incident customer follow-up from CEO Ariel Nazarova",
         "Full public post-mortem at halcyon.security/postmortems/<case-id> with timeline, technical details, remediation",
         "Internal blameless review led by Elena Kowalski with IR team within 14 days",
         "Update this playbook based on after-action findings — Sonia tracks in Asana project 'crisis-playbook-iter'",
         "Media drill to test the updated playbook every quarter (next: 2026-06-14)"])

    path = f"{OUT}/halcyon_crisis_comms_playbook.pptx"
    p.save(path)
    return path


if __name__ == "__main__":
    paths = [
        docx_meridian_claim_triage(),
        docx_kelvingrove_onboarding(),
        docx_zephyr_dsar_compliance(),
        pptx_orbital_enterprise_upsell(),
        pptx_lumen_trial047(),
        pptx_halcyon_crisis_comms(),
    ]
    for p in paths:
        size = os.path.getsize(p)
        print(f"  {p} ({size} bytes)")
    print(f"\n{len(paths)} files in {OUT}")
